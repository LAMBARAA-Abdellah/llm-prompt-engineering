#!/usr/bin/env python3
"""
Generate PowerPoint presentation for the Prompt Engineering project.

Creates a 15-20 slide professional presentation covering all project aspects.
Run: python docs/generate_presentation.py
"""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)


def generate_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = [
        ("Prompt Engineering with Python using LLMs",
         "Master BDCC — Université Mohammed VI Polytechnique\nAcademic Project 2026"),
        ("Agenda", [
            "Introduction to Prompt Engineering",
            "Large Language Models Overview",
            "Project Architecture",
            "LLM Providers: OpenAI, Ollama, Groq",
            "Prompt Engineering Techniques",
            "Use Cases: Sentiment Analysis & Spam Detection",
            "Image Prompting (Text-to-Image, Image-to-Image)",
            "Evaluation Methodology",
            "Results & Comparison",
            "Conclusion & Future Work",
        ]),
        ("Introduction", [
            "Prompt Engineering is the art and science of crafting inputs to LLMs",
            "Goal: Maximize model performance without fine-tuning",
            "Critical skill for NLP practitioners and AI engineers",
            "This project implements and compares 8 prompting techniques",
            "Evaluated on real-world NLP tasks with measurable metrics",
        ]),
        ("What are Large Language Models?", [
            "Neural networks trained on vast text corpora",
            "Generate human-like text based on input prompts",
            "Examples: GPT-4, Llama 3, Mixtral, Gemma",
            "Accessed via APIs (OpenAI, Groq) or locally (Ollama)",
            "Performance heavily depends on prompt quality",
        ]),
        ("Project Objectives", [
            "Develop a modular Python application for prompt engineering",
            "Support multiple LLM providers (OpenAI, Ollama, Groq)",
            "Implement 8 prompting techniques with code examples",
            "Evaluate techniques on Sentiment Analysis and Spam Detection",
            "Compare accuracy, speed, and cost across techniques",
            "Demonstrate text-to-image and image-to-image prompting",
        ]),
        ("System Architecture", [
            "Clean layered architecture with separation of concerns",
            "config/ — Environment-based settings and model registry",
            "llms/ — Provider-specific clients with retry logic",
            "prompts/ — One module per prompting technique",
            "use_cases/ — Business logic orchestrating workflows",
            "evaluation/ — Metrics, comparison, cost estimation",
        ]),
        ("LLM Providers", [
            "OpenAI: GPT-4o-mini, DALL-E 3 — Cloud API with image generation",
            "Ollama: Llama3, Mistral, Gemma — Local inference, zero cost",
            "Groq: Llama 3.3 70B, Mixtral — Fast cloud inference",
            "All clients share BaseLLMClient interface",
            "Automatic retry, timing, and token tracking",
        ]),
        ("Prompt Engineering Techniques (1/2)", [
            "Simple Prompt: Direct instruction, no context",
            "Zero-Shot: Task description without examples",
            "Few-Shot: 2-5 labeled examples before input",
            "Chain of Thought: Step-by-step reasoning",
        ]),
        ("Prompt Engineering Techniques (2/2)", [
            "Self-Consistency: Multiple samples + majority voting",
            "Role Prompt: Expert persona via system message",
            "Step-Back: Derive principles, then apply",
            "Tree of Thoughts: Multi-branch reasoning + selection",
        ]),
        ("Use Case 1: Sentiment Analysis", [
            "Dataset: IMDB Movie Reviews (positive/negative)",
            "20 samples with gold labels for evaluation",
            "6 prompting techniques compared",
            "Metrics: Accuracy, Precision, Recall, F1 Score",
            "EDA visualizations: label distribution, text length",
        ]),
        ("Use Case 2: Spam Detection", [
            "Dataset: SMS Spam Collection (ham/spam)",
            "Same evaluation pipeline as sentiment analysis",
            "5 prompting techniques tested",
            "Gold examples for few-shot learning",
            "Confusion matrices per technique",
        ]),
        ("Image Prompting", [
            "Text-to-Image: DALL-E 3 generates images from descriptions",
            "Image-to-Image: DALL-E 2 edits existing images",
            "Automatic saving with prompt metadata",
            "Demonstrates multimodal prompt engineering",
        ]),
        ("Gold Examples", [
            "Manually labeled high-quality reference data",
            "Essential for benchmarking and reproducibility",
            "Used directly in few-shot prompting",
            "8 gold examples per task (sentiment + spam)",
            "Each includes rationale for the label",
        ]),
        ("Evaluation Methodology", [
            "Metrics: Accuracy, Precision, Recall, F1 Score",
            "Confusion Matrix per technique",
            "Execution Time tracking per prediction",
            "Token Usage and Cost Estimation",
            "Comparison charts: accuracy vs speed vs cost",
        ]),
        ("Expected Results", [
            "Few-Shot and Role Prompt typically achieve highest accuracy",
            "Zero-Shot provides good baseline with lowest cost",
            "Chain of Thought improves interpretability",
            "Self-Consistency reduces variance at higher cost",
            "Ollama offers free local inference vs cloud API costs",
        ]),
        ("Technology Stack", [
            "Python 3.12 with type hints and PEP8 compliance",
            "openai, groq, ollama SDKs",
            "pandas, scikit-learn for data and metrics",
            "matplotlib, seaborn for visualization",
            "Rich for beautiful console output",
            "pytest for testing, Jupyter for exploration",
        ]),
        ("Conclusion", [
            "Comprehensive prompt engineering framework implemented",
            "8 techniques compared on 2 real NLP tasks",
            "Modular, extensible, production-ready architecture",
            "Multimodal capabilities demonstrated with DALL-E",
            "Clear trade-offs between accuracy, speed, and cost",
        ]),
        ("Future Work", [
            "Fine-tuning comparison with prompt engineering",
            "Retrieval-Augmented Generation (RAG) integration",
            "Automated prompt optimization (DSPy, AutoPrompt)",
            "Multi-language support for datasets",
            "Real-time streaming evaluation dashboard",
            "Integration with LangChain/LlamaIndex frameworks",
        ]),
        ("Thank You", [
            "Questions & Discussion",
            "GitHub: PromptEngineeringProject",
            "Author: Master BDCC Student",
            "Université Mohammed VI Polytechnique — 2026",
        ]),
    ]

    add_title_slide(prs, slides_data[0][0], slides_data[0][1])
    for title, bullets in slides_data[1:]:
        if title == "Thank You":
            add_title_slide(prs, title, "\n".join(bullets))
        else:
            add_content_slide(prs, title, bullets)

    output_path = PROJECT_ROOT / "outputs" / "Prompt_Engineering_Presentation.pptx"
    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_presentation()
